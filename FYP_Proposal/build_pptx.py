"""Build the Sift FYP Proposal slide deck — 7 slides, ~3:45 narration."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# palette
BG       = RGBColor(0xFF, 0xFF, 0xFF)
INK      = RGBColor(0x10, 0x1A, 0x2E)
MUTED    = RGBColor(0x55, 0x60, 0x72)
ACCENT   = RGBColor(0x00, 0x8E, 0x86)
HAIRLINE = RGBColor(0xE3, 0xE6, 0xEC)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank_layout = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(blank_layout)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def add_header(slide, eyebrow, title):
    add_text(slide, Inches(0.7), Inches(0.6), Inches(8), Inches(0.3),
             eyebrow, size=11, bold=True, color=ACCENT)
    add_text(slide, Inches(0.7), Inches(0.95), Inches(12), Inches(0.9),
             title, size=32, bold=True, color=INK)


def add_footer(slide, page):
    add_rect(slide, Inches(0.7), Inches(7.05), Inches(11.95), Emu(6350),
             fill=HAIRLINE)
    add_text(slide, Inches(0.7), Inches(7.15), Inches(8), Inches(0.3),
             "Sift  ·  CM3070 Final Project  ·  Template 4.2 — Financial Advisor Bot",
             size=10, color=MUTED)
    add_text(slide, Inches(11.6), Inches(7.15), Inches(1.0), Inches(0.3),
             str(page), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# SLIDE 1 — Title
# ---------------------------------------------------------------------------
s = add_slide()

add_text(s, Inches(0.9), Inches(0.85), Inches(2), Inches(0.45),
         "PROJECT PROPOSAL", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.9), Inches(2.5), Inches(11), Inches(1.5),
         "Sift", size=96, bold=True, color=INK)
add_text(s, Inches(0.9), Inches(4.1), Inches(11), Inches(0.6),
         "An AI tool that explains company earnings to non-expert investors.",
         size=24, color=MUTED)
add_text(s, Inches(0.9), Inches(6.4), Inches(11), Inches(0.4),
         "CM3070 Computer Science Final Project  ·  Template 4.2 — Financial Advisor Bot",
         size=13, bold=True, color=INK)
add_text(s, Inches(0.9), Inches(6.75), Inches(11), Inches(0.4),
         "Leroy Ng  ·  University of London",
         size=12, color=MUTED)

# ---------------------------------------------------------------------------
# SLIDE 2 — Motivation
# ---------------------------------------------------------------------------
s = add_slide()
add_header(s, "MOTIVATION",
           "Most everyday investors can't parse what companies actually publish")

col_w = Inches(5.8); col_y = Inches(2.4)

add_text(s, Inches(0.7), col_y, col_w, Inches(0.35),
         "WHAT COMPANIES PUBLISH", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.7), col_y + Inches(0.4), col_w, Inches(0.6),
         "Dense filings · hour-long calls", size=22, bold=True, color=INK)
add_text(s, Inches(0.7), col_y + Inches(1.1), col_w, Inches(2.4),
         "Formal 8-K filings, multi-page press releases, and earnings call "
         "transcripts written for professional analysts. Heavy on technical "
         "terms (guidance, margins, EBITDA, free cash flow etc)",
         size=15, color=MUTED)

add_text(s, Inches(6.85), col_y, col_w, Inches(0.35),
         "WHAT MOST RETAIL USERS SEE", size=11, bold=True, color=ACCENT)
add_text(s, Inches(6.85), col_y + Inches(0.4), col_w, Inches(0.6),
         "A date · a single EPS number", size=22, bold=True, color=INK)
add_text(s, Inches(6.85), col_y + Inches(1.1), col_w, Inches(2.4),
         "Yahoo Finance and broker apps show the release date and the "
         "headline earnings figure. No context, no explanation of what any "
         "of it means.",
         size=15, color=MUTED)

add_rect(s, Inches(0.7), Inches(6.5), Inches(11.95), Emu(6350), fill=HAIRLINE)
add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
         "The project asks whether current AI can turn dense filings into plain-English explanations users can learn from.",
         size=14, color=INK, bold=True)

add_footer(s, 2)

# ---------------------------------------------------------------------------
# SLIDE 3 — Key terms
# ---------------------------------------------------------------------------
s = add_slide()
add_header(s, "KEY TERMS", "The terms used in this project.")

terms = [
    ("SEC EDGAR",
     "The US government's free, public database of company filings.  "
     "Used in this project as the source of raw data."),
    ("8-K filing",
     "A short report companies must publish when something important happens, "
     "including earnings results.  Used in this project as the signal to start "
     "analysing a release."),
    ("Earnings call transcript",
     "A written record of the executives quarterly briefing to investors.  "
     "Used in this project to detect shifts in tone and outlook."),
    ("Consensus estimate",
     "The average prediction analysts make for a company's earnings.  "
     "Used in this project as context to help the user see whether the "
     "actual result was above or below expectations."),
]
ry = Inches(2.25); rh = Inches(1.05); rg = Inches(0.12)
for i, (term, body) in enumerate(terms):
    y = ry + (rh + rg) * i
    add_text(s, Inches(0.7), y, Inches(3.3), rh,
             term, size=18, bold=True, color=INK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.2), y, Inches(8.5), rh,
             body, size=13, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(terms) - 1:
        add_rect(s, Inches(0.7), y + rh + rg/2 - Emu(3175),
                 Inches(11.95), Emu(6350), fill=HAIRLINE)

add_footer(s, 3)

# ---------------------------------------------------------------------------
# SLIDE 4 — Related work
# ---------------------------------------------------------------------------
s = add_slide()
add_header(s, "RELATED WORK", "Four works that ground the design")

works = [
    ("ACADEMIC", "FinBERT  ·  Araci, 2019",
     "BERT fine-tuned for financial text.",
     "GPT-4o with careful prompting now beats it by ~10% (MDPI 2025). Justifies using a general-purpose LLM over a specialised one."),
    ("ACADEMIC", "Financial literacy  ·  Lusardi & Mitchell, 2014",
     "Retail investors lack the vocabulary and context to parse financial information.",
     "Establishes the underlying gap the project aims to address with plain-English explanations."),
    ("COMMERCIAL", "AlphaSense",
     "Professional-grade earnings synthesis.",
     "Proves the synthesis task is achievable, but enterprise-only and not built for non-experts."),
    ("COMMERCIAL", "Yahoo Finance · Broker apps",
     "Retail-facing surfaces for earnings data.",
     "Show the date and a headline number, but offer no explanation, context, or education."),
]
cw = Inches(5.95); ch = Inches(2.15); gx = Inches(0.15); gy = Inches(0.15)
ox = Inches(0.7); oy = Inches(2.25)
for i, (cat, name, headline, critique) in enumerate(works):
    col, row = i % 2, i // 2
    x = ox + (cw + gx) * col
    y = oy + (ch + gy) * row
    add_text(s, x, y, cw, Inches(0.32),
             cat, size=10, bold=True, color=ACCENT)
    add_text(s, x, y + Inches(0.32), cw, Inches(0.4),
             name, size=13, bold=True, color=INK)
    add_text(s, x, y + Inches(0.78), cw, Inches(0.5),
             headline, size=15, bold=True, color=INK)
    add_text(s, x, y + Inches(1.35), cw, Inches(0.8),
             critique, size=12, color=MUTED)
add_rect(s, ox, oy + ch + Inches(0.02), cw * 2 + gx, Emu(6350), fill=HAIRLINE)

add_footer(s, 4)

# ---------------------------------------------------------------------------
# SLIDE 5 — Synthesis
# ---------------------------------------------------------------------------
s = add_slide()
add_header(s, "SYNTHESIS", "What the related work shows — and what's still open.")

rows = [
    ("Established",
     "Transformer-based NLP works well on financial text."),
    ("Established",
     "Retail investors face a documented financial-literacy gap."),
    ("Established",
     "High-end synthesis tools exist for institutional users."),
    ("Open",
     "An AI tool that explains earnings releases / stock analyses in plain language to non-expert users."),
]
ry = Inches(2.3); rh = Inches(0.85); rg = Inches(0.15)
for i, (k, body) in enumerate(rows):
    y = ry + (rh + rg) * i
    is_open = k == "Open"
    add_text(s, Inches(0.7), y, Inches(2.0), rh,
             k.upper(), size=12, bold=True,
             color=INK if is_open else MUTED,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.8), y, Inches(9.85), rh,
             body, size=14, bold=is_open,
             color=INK if is_open else MUTED,
             anchor=MSO_ANCHOR.MIDDLE)
    if i < len(rows) - 1:
        add_rect(s, Inches(0.7), y + rh + rg/2 - Emu(3175),
                 Inches(11.95), Emu(6350), fill=HAIRLINE)

add_footer(s, 5)

# ---------------------------------------------------------------------------
# SLIDE 6 — Approach
# ---------------------------------------------------------------------------
s = add_slide()
add_header(s, "APPROACH", "What gets built, and how.")

# Top: 4 components in a compact row
comp_y = Inches(2.15)
add_text(s, Inches(0.7), comp_y, Inches(2), Inches(0.3),
         "COMPONENTS", size=11, bold=True, color=ACCENT)
comps = [
    ("Pre-earnings briefing",
     "LLM writes a plain-English summary of what to expect from the upcoming release."),
    ("Filing explainer",
     "When an 8-K hits, the LLM produces a plain-English explanation of what was reported."),
    ("Transcript breakdown",
     "LLM summarises the earnings call — key points, tone, and any change in outlook."),
    ("Ask Sift",
     "User can ask follow-up questions about any term, figure, or section they don't understand."),
]
cw = Inches(2.95); cg = Inches(0.13)
cy = comp_y + Inches(0.4)
for i, (t, b) in enumerate(comps):
    x = Inches(0.7) + (cw + cg) * i
    add_text(s, x, cy, cw, Inches(0.45),
             t, size=13, bold=True, color=INK)
    add_text(s, x, cy + Inches(0.45), cw, Inches(1.6),
             b, size=11, color=MUTED)

# Divider
add_rect(s, Inches(0.7), Inches(4.8), Inches(11.95), Emu(6350), fill=HAIRLINE)

# Bottom: stack + evaluation in two columns
btm_y = Inches(5.0)
add_text(s, Inches(0.7), btm_y, Inches(3), Inches(0.3),
         "TECHNICAL STACK", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.7), btm_y + Inches(0.4), Inches(5.8), Inches(1.6),
         "Python workers on Modal for ingestion and LLM orchestration.\n"
         "Supabase for storage and user accounts.\n"
         "React Native mobile app on iOS and Android.\n"
         "AI core: a large language model (Claude or GPT) handles all "
         "summarisation, explanation, and Q&A.",
         size=12, color=MUTED)

add_text(s, Inches(6.85), btm_y, Inches(3), Inches(0.3),
         "EVALUATION", size=11, bold=True, color=ACCENT)
add_text(s, Inches(6.85), btm_y + Inches(0.4), Inches(5.8), Inches(1.6),
         "Manual review of LLM outputs against source filings for accuracy.\n"
         "Small user study on whether the explanations actually aid understanding.",
         size=12, color=MUTED)

add_footer(s, 6)

# ---------------------------------------------------------------------------
# SLIDE 7 — Deliverables
# ---------------------------------------------------------------------------
s = add_slide()
add_header(s, "DELIVERABLES", "What the submission will contain.")

items = [
    ("Working pipeline",
     "End-to-end ingestion and LLM-based explanation, running on 50 selected tickers."),
    ("Evaluation report",
     "Accuracy review against source filings, readability scoring, and user-study results."),
    ("Mobile application",
     "A React Native app showing the explanations and Q&A on iOS and Android."),
    ("Documentation",
     "Architecture notes, prompt design, data sources, and ethics positioning."),
]
ry = Inches(2.3); rh = Inches(0.95); rg = Inches(0.15)
for i, (title, body) in enumerate(items):
    y = ry + (rh + rg) * i
    add_text(s, Inches(0.7), y, Inches(0.5), rh,
             "·", size=24, bold=True, color=ACCENT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.2), y, Inches(4.0), rh,
             title, size=18, bold=True, color=INK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.4), y, Inches(7.2), rh,
             body, size=14, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(items) - 1:
        add_rect(s, Inches(0.7), y + rh + rg/2 - Emu(3175),
                 Inches(11.95), Emu(6350), fill=HAIRLINE)

add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
         "Framed as a research and education project — not financial advice.",
         size=12, color=MUTED)

add_footer(s, 7)


out = r"C:\Users\leheh\.Projects\Sift\FYP_Proposal\Sift_Proposal.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
